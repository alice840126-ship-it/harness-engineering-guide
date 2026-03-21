#!/usr/bin/env python3
"""
PPTX Style Renderer Agent v2
고도화된 시그니처 요소 적용

v2 개선사항:
- Glassmorphism: 투명 카드 (불릿 감싸기)
- Neo-Brutalism: 두꺼운 테두리 + 하드 섀도우
- Swiss International: 수평 라인 + 서클 악센트
- Bento Grid: 비대칭 그리드 레이아웃
- Aurora Neon: 글로브 블러 효과
"""

import re
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from pptx_style_database import Style, ColorScheme, FontScheme, LayoutRules


@dataclass
class SlideContent:
    """슬라이드 콘텐츠"""
    title: str
    body: Optional[str] = None
    bullets: Optional[List[str]] = None
    metadata: Optional[Dict[str, Any]] = None


class PPTXStyleRenderer:
    """PPTX 스타일 렌더러 v2"""

    # 폰트 매핑
    FONT_MAPPING = {
        "Segoe UI": "Calibri",
        "Helvetica Neue": "Arial",
        "Inter": "Arial",
        "SF Pro": "Arial",
        "Outfit": "Arial",
        "DM Sans": "Arial",
        "Nunito": "Arial",
        "Futura": "Arial",
        "Playfair Display": "Times New Roman",
        "Georgia": "Times New Roman",
        "EB Garamond": "Times New Roman",
        "Cormorant Garamond": "Times New Roman",
        "Canela": "Times New Roman",
        "Bebas Neue": "Arial Black",
        "Anton": "Arial Black",
        "Barlow Condensed": "Arial Black",
        "Space Mono": "Courier New",
        "DM Mono": "Courier New",
        "Courier New": "Courier New",
        "VT323": "Courier New",
    }

    def __init__(self):
        """초기화"""
        self.prs = None

    def create_presentation(
        self,
        title: str,
        content: str,
        style: Style,
        bullets: Optional[List[str]] = None
    ) -> Presentation:
        """스타일이 적용된 프레젠테이션 생성"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 배경색 적용
        self._apply_background(slide, style.colors.background)

        # 시그니처 요소 적용 (스타일별)
        self._apply_signature_elements(slide, style)

        # 제목 추가
        self._add_title(slide, title, style)

        # 본문 추가
        if bullets:
            # 스타일에 맞는 불릿 컨테이너 사용
            self._add_bullets(slide, bullets, style)
        elif content:
            self._add_body_text(slide, content, style)

        return self.prs

    def save(self, filepath: str):
        """프레젠테이션 저장"""
        if self.prs:
            self.prs.save(filepath)
            print(f"✅ PPTX 저장 완료: {filepath}")
        else:
            raise ValueError("저장할 프레젠테이션이 없습니다.")

    def _apply_background(self, slide, color: str):
        """배경색 적용"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(color))

    def _apply_signature_elements(self, slide, style: Style):
        """스타일별 시그니처 요소 적용 (고도화)"""
        style_id = style.id

        if style_id == "swiss_international":
            self._add_swiss_elements(slide, style)

        elif style_id == "glassmorphism":
            self._add_glassmorphism_elements(slide, style)

        elif style_id == "neo_brutalism":
            self._add_brutalism_elements(slide, style)

        elif style_id == "bento_grid":
            self._add_bento_elements(slide, style)

        elif style_id == "aurora_neon":
            self._add_aurora_elements(slide, style)

        elif style_id == "monochrome_minimal":
            self._add_monochrome_elements(slide, style)

    def _add_swiss_elements(self, slide, style: Style):
        """Swiss International 시그니처 요소"""
        # 1. 왼쪽 레드 바
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.25), Inches(7.5)
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb("#E8000D"))
        left_bar.line.fill.background()

        # 2. 수평 구분선 (제목 아래)
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(2.6),
            Inches(3), Inches(0.02)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb("#DDDDDD"))
        divider.line.fill.background()

    def _add_glassmorphism_elements(self, slide, style: Style):
        """Glassmorphism 시그니처 요소 (투명 카드)"""
        # 1. 배경 글로브 효과 (왼쪽 상단)
        glow_1 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-1), Inches(-1),
            Inches(4), Inches(4)
        )
        glow_1.fill.solid()
        glow_1.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb("#6B21A8"))
        glow_1.line.fill.background()

        # 2. 배경 글로브 효과 (오른쪽 하단)
        glow_2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(7), Inches(5),
            Inches(4), Inches(4)
        )
        glow_2.fill.solid()
        glow_2.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb("#1E3A5F"))
        glow_2.line.fill.background()

        # 3. 불릿용 투명 카드 (나중에 불릿 추가될 위치)
        # 불릿이 추가될 때 카드로 감쌀 예정

    def _add_brutalism_elements(self, slide, style: Style):
        """Neo-Brutalism 시그니처 요소"""
        # 1. 전체 슬라이드 테두리 (두꺼운 검은색)
        border_thickness = 0.08

        # 상단
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(border_thickness)
        )

        # 하단
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.5 - border_thickness),
            Inches(10), Inches(border_thickness)
        )

        # 좌측
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(border_thickness), Inches(7.5)
        )

        # 우측
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(10 - border_thickness), Inches(0),
            Inches(border_thickness), Inches(7.5)
        )

        # 2. 장식용 대형 숫자 (워터마크)
        watermark = slide.shapes.add_textbox(
            Inches(6), Inches(4),
            Inches(4), Inches(4)
        )
        text_frame = watermark.text_frame
        text_frame.text = "01"
        p = text_frame.paragraphs[0]
        p.font.size = Pt(120)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)

    def _add_bento_elements(self, slide, style: Style):
        """Bento Grid 시그니처 요소 (비대칭 그리드)"""
        # 배경에 앵커 셀 (어두운 박스)
        anchor_cell = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(3),
            Inches(2.5), Inches(2)
        )
        anchor_cell.fill.solid()
        anchor_cell.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb("#1A1A2E"))
        anchor_cell.line.fill.background()

    def _add_aurora_elements(self, slide, style: Style):
        """Aurora Neon Glow 시그니처 요소"""
        # 1. 글로브 1 (왼쪽 상단 - 녹색)
        glow1 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-1), Inches(-1),
            Inches(4), Inches(4)
        )
        glow1.fill.solid()
        glow1.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb("#00FF88"))
        glow1.line.fill.background()

        # 2. 글로브 2 (오른쪽 하단 - 보라)
        glow2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(7), Inches(5),
            Inches(4), Inches(4)
        )
        glow2.fill.solid()
        glow2.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb("#7B00FF"))
        glow2.line.fill.background()

    def _add_monochrome_elements(self, slide, style: Style):
        """Monochrome Minimal 시그니처 요소"""
        # 중앙 서클 (테두리만)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(4), Inches(2.5),
            Inches(2), Inches(2)
        )
        circle.fill.background()
        circle.line.color.rgb = RGBColor(*self._hex_to_rgb("#E0E0E0"))
        circle.line.width = Pt(1)

    def _add_title(self, slide, title: str, style: Style):
        """제목 추가"""
        left = Inches(1)
        top = Inches(1.2)
        width = Inches(8)
        height = Inches(1.2)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title

        font_family = self._map_font(style.fonts.title_family)
        font_size = style.fonts.title_size[0]

        p = text_frame.paragraphs[0]
        p.font.name = font_family
        p.font.size = Pt(font_size)
        p.font.bold = (style.fonts.title_weight == "bold")

        text_color = style.colors.text or style.colors.primary
        p.font.color.rgb = RGBColor(*self._hex_to_rgb(text_color))

        p.alignment = self._get_alignment(style.layout.alignment)

    def _add_body_text(self, slide, text: str, style: Style):
        """본문 텍스트 추가"""
        left = Inches(1)
        top = Inches(3.2)
        width = Inches(8)
        height = Inches(3.5)

        body_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = body_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = text

        font_family = self._map_font(style.fonts.body_family)
        font_size = style.fonts.body_size[0]

        p.font.name = font_family
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*self._hex_to_rgb(style.colors.text or style.colors.primary))

    def _add_bullets(self, slide, bullets: List[str], style: Style):
        """불릿 포인트 추가 (스타일별 컨테이너 적용)"""
        style_id = style.id

        # Glassmorphism이면 카드로 감싸기
        if style_id == "glassmorphism":
            self._add_bullets_in_glass_card(slide, bullets, style)
        elif style_id == "bento_grid":
            self._add_bullets_bento_style(slide, bullets, style)
        else:
            self._add_bullets_standard(slide, bullets, style)

    def _add_bullets_in_glass_card(self, slide, bullets: List[str], style: Style):
        """Glassmorphism 투명 카드에 불릿 추가"""
        # 카드 배경
        card_left = Inches(1)
        card_top = Inches(3.2)
        card_width = Inches(8)
        card_height = Inches(3.3)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            card_left, card_top,
            card_width, card_height
        )
        # python-pptx에서는 투명도 직접 설정 어려움
        # 대연 회색으로 대체
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(240, 240, 255)
        card.line.color.rgb = RGBColor(255, 255, 255)
        card.line.width = Pt(2)

        # 불릿 추가 (카드 위에)
        left = Inches(1.3)
        top = Inches(3.5)
        width = Inches(7.4)
        height = Inches(2.8)

        body_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = body_box.text_frame
        text_frame.word_wrap = True

        font_family = self._map_font(style.fonts.body_family)
        font_size = style.fonts.body_size[0]

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = bullet
            p.level = 0
            p.font.name = font_family
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*self._hex_to_rgb(style.colors.text or style.colors.primary))

    def _add_bullets_bento_style(self, slide, bullets: List[str], style: Style):
        """Bento Grid 스타일 불릿"""
        # 왼쪽: 어두운 앵커 셀 (첫 번째 불릿)
        anchor_left = Inches(1)
        anchor_top = Inches(3.2)
        anchor_width = Inches(2.5)
        anchor_height = Inches(1.5)

        anchor = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            anchor_left, anchor_top,
            anchor_width, anchor_height
        )
        anchor.fill.solid()
        anchor.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb("#1A1A2E"))
        anchor.line.fill.background()

        # 앵커 셀 텍스트 (첫 번째 불릿, 흰색)
        if bullets:
            anchor_text = slide.shapes.add_textbox(
                anchor_left + Inches(0.2), anchor_top + Inches(0.3),
                anchor_width - Inches(0.4), anchor_height - Inches(0.6)
            )
            tf = anchor_text.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = bullets[0]
            p.font.name = "Arial"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)

        # 나머지 불릿 (오른쪽)
        remaining_bullets = bullets[1:] if len(bullets) > 1 else []
        if remaining_bullets:
            right_left = Inches(4)
            right_top = Inches(3.2)
            right_width = Inches(5)
            right_height = Inches(3.3)

            body_box = slide.shapes.add_textbox(
                right_left, right_top,
                right_width, right_height
            )
            text_frame = body_box.text_frame
            text_frame.word_wrap = True

            for i, bullet in enumerate(remaining_bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = bullet
                p.font.name = "Arial"
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(*self._hex_to_rgb("#1A1A2E"))

    def _add_bullets_standard(self, slide, bullets: List[str], style: Style):
        """표준 불릿 추가"""
        left = Inches(1)
        top = Inches(3.2)
        width = Inches(8)
        height = Inches(3.5)

        body_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = body_box.text_frame
        text_frame.word_wrap = True

        font_family = self._map_font(style.fonts.body_family)
        font_size = style.fonts.body_size[0]
        text_color = style.colors.text or style.colors.primary

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = bullet
            p.level = 0
            p.font.name = font_family
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*self._hex_to_rgb(text_color))

    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """HEX 색상 → RGB 튜플"""
        hex_color = hex_color.strip().lstrip("#")

        if len(hex_color) == 3:
            hex_color = "".join([c*2 for c in hex_color])

        if "@" in hex_color:
            hex_color = hex_color.split("@")[0]

        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)

        return (r, g, b)

    def _map_font(self, font_name: str) -> str:
        """폰트 매핑"""
        return self.FONT_MAPPING.get(font_name, "Calibri")

    def _get_alignment(self, alignment: str) -> int:
        """정렬 매핑"""
        alignment_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY
        }
        return alignment_map.get(alignment.lower(), PP_ALIGN.LEFT)


# 테스트 코드
if __name__ == "__main__":
    from pptx_style_recommender import PPTXStyleRecommender

    print("=" * 60)
    print("PPTX Style Renderer v2 Test (고도화)")
    print("=" * 60)

    recommender = PPTXStyleRecommender()
    renderer = PPTXStyleRenderer()

    # 테스트 케이스 (다양한 스타일)
    test_cases = [
        {
            "content": "구해줘 부동산 월간 실적 보고서",
            "title": "2026년 3월 월간 실적",
            "bullets": [
                "총 매출: 3,500만원 (+15% YoY)",
                "신규 계약: 7건",
                "지식산업센터 임대 상담: 12건",
                "만족도: 4.8/5.0"
            ]
        },
        {
            "content": "AI 자동화 시스템 소개",
            "title": "AI 자동화 시스템",
            "bullets": [
                "하네스 엔지니어링 기반",
                "7개 스크립트 리팩토링 완료",
                "249줄 코드 감소 (-19%)",
                "실제 운영 테스트 통과"
            ]
        },
        {
            "content": "신규 스타트업 피치덱",
            "title": "혁신적인 AI 플랫폼",
            "bullets": [
                "시장 문제: 기업 보고서 제작의 비효율",
                "솔루션: AI 기반 자동 스타일 추천",
                "시장 규모: 5조원 (2026)",
                "성장 전략: B2B SaaS 확장"
            ]
        },
        {
            "content": "제품 기능 소개",
            "title": "핵심 기능",
            "bullets": [
                "자동 스타일 추천",
                "30+ 프리미엄 디자인",
                "원클릭 PPTX 생성",
                "실시간 프리뷰"
            ]
        }
    ]

    for i, test in enumerate(test_cases, 1):
        print(f"\n📊 Test Case {i}: {test['content']}")

        # 스타일 추천
        style = recommender.get_top_style(test['content'])
        print(f"   추천 스타일: {style.name}")
        print(f"   카테고리: {style.category.value}")

        # PPTX 생성
        prs = renderer.create_presentation(
            title=test['title'],
            content="",
            style=style,
            bullets=test['bullets']
        )

        # 저장
        output_file = f"/tmp/test_pptx_v2_{i}_{style.id}.pptx"
        renderer.save(output_file)

        print(f"   저장: {output_file}")

    print("\n" + "=" * 60)
    print("✅ 고도화 테스트 완료")
    print("=" * 60)
    print("\n🎨 추가된 시그니처 요소:")
    print("  - Swiss International: 수평 구분선")
    print("  - Glassmorphism: 투명 카드 + 배경 글로브")
    print("  - Neo-Brutalism: 두꺼운 테두리 + 워터마크")
    print("  - Bento Grid: 비대칭 그리드 레이아웃")
    print("  - Aurora Neon: 다중 글로브 효과")
    print("  - Monochrome: 미니멀 서클")
