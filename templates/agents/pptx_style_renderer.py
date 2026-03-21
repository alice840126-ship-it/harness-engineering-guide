#!/usr/bin/env python3
"""
PPTX Style Renderer Agent
스타일 → python-pptx로 렌더링

Usage:
    from pptx_style_renderer import PPTXStyleRenderer
    from pptx_style_recommender import PPTXStyleRecommender

    recommender = PPTXStyleRecommender()
    renderer = PPTXStyleRenderer()

    # 자동 스타일 추천 + PPTX 생성
    style = recommender.get_top_style("AI 자동화 보고서")
    pptx = renderer.create_presentation(
        title="AI 자동화 시스템",
        content="내용...",
        style=style
    )
    pptx.save("output.pptx")
"""

import re
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from pptx_style_database import Style, ColorScheme, FontScheme, LayoutRules


@dataclass
class SlideContent:
    """슬라이드 콘텐츠"""
    title: str
    body: Optional[str] = None
    bullets: Optional[List[str]] = None
    metadata: Optional[Dict[str, Any]] = None


class PPTXStyleRenderer:
    """PPTX 스타일 렌더러"""

    # 폰트 매핑 (스타일 폰트 → 시스템 폰트)
    FONT_MAPPING = {
        # Sans-serif
        "Segoe UI": "Calibri",
        "Helvetica Neue": "Arial",
        "Inter": "Arial",
        "SF Pro": "Arial",
        "Outfit": "Arial",
        "DM Sans": "Arial",
        "Nunito": "Arial",
        "Futura": "Arial",

        # Serif
        "Playfair Display": "Times New Roman",
        "Georgia": "Times New Roman",
        "EB Garamond": "Times New Roman",
        "Cormorant Garamond": "Times New Roman",
        "Canela": "Times New Roman",

        # Display/Headline
        "Bebas Neue": "Arial Black",
        "Anton": "Arial Black",
        "Barlow Condensed": "Arial Black",

        # Monospace
        "Space Mono": "Courier New",
        "DM Mono": "Courier New",
        "Courier New": "Courier New",
        "VT323": "Courier New",
    }

    def __init__(self):
        """초기화"""
        self.prs = None

    def create_presentation(
        self,
        title: str,
        content: str,
        style: Style,
        bullets: Optional[List[str]] = None
    ) -> Presentation:
        """스타일이 적용된 프레젠테이션 생성

        Args:
            title: 슬라이드 제목
            content: 본문 내용
            style: 적용할 스타일
            bullets: 불릿 포인트 리스트 (선택)

        Returns:
            Presentation 객체
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        # 슬라이드 추가 (빈 레이아웃)
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # 배경색 적용
        self._apply_background(slide, style.colors.background)

        # 시그니처 요소 적용 (배경 장식)
        self._apply_signature_elements(slide, style)

        # 제목 추가
        self._add_title(slide, title, style)

        # 본문 추가 (불릿 또는 텍스트)
        if bullets:
            self._add_bullets(slide, bullets, style)
        elif content:
            self._add_body_text(slide, content, style)

        return self.prs

    def create_multi_slide_presentation(
        self,
        slides_data: List[SlideContent],
        style: Style
    ) -> Presentation:
        """여러 슬라이드 프레젠테이션 생성

        Args:
            slides_data: 슬라이드 데이터 리스트
            style: 적용할 스타일

        Returns:
            Presentation 객체
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        for slide_data in slides_data:
            self.create_presentation(
                title=slide_data.title,
                content=slide_data.body,
                style=style,
                bullets=slide_data.bullets
            )

        return self.prs

    def save(self, filepath: str):
        """프레젠테이션 저장"""
        if self.prs:
            self.prs.save(filepath)
            print(f"✅ PPTX 저장 완료: {filepath}")
        else:
            raise ValueError("저장할 프레젠테이션이 없습니다.")

    def _apply_background(self, slide, color: str):
        """배경색 적용"""
        background = slide.background
        fill = background.fill
        fill.solid()

        # HEX → RGB
        rgb = self._hex_to_rgb(color)
        fill.fore_color.rgb = RGBColor(*rgb)

    def _add_title(self, slide, title: str, style: Style):
        """제목 추가"""
        # 제목 텍스트박스
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(1)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title

        # 폰트 설정
        font_family = self._map_font(style.fonts.title_family)
        font_size = style.fonts.title_size[0]  # 최소 크기 사용

        p = text_frame.paragraphs[0]
        p.font.name = font_family
        p.font.size = Pt(font_size)
        p.font.bold = (style.fonts.title_weight == "bold")

        # 색상
        text_color = style.colors.text or style.colors.primary
        p.font.color.rgb = RGBColor(*self._hex_to_rgb(text_color))

        # 정렬
        p.alignment = self._get_alignment(style.layout.alignment)

    def _add_body_text(self, slide, text: str, style: Style):
        """본문 텍스트 추가"""
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(3.5)

        body_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = body_box.text_frame
        text_frame.word_wrap = True

        # 텍스트 추가
        p = text_frame.paragraphs[0]
        p.text = text

        # 폰트 설정
        font_family = self._map_font(style.fonts.body_family)
        font_size = style.fonts.body_size[0]

        p.font.name = font_family
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*self._hex_to_rgb(style.colors.text or style.colors.primary))

    def _add_bullets(self, slide, bullets: List[str], style: Style):
        """불릿 포인트 추가"""
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(3.5)

        body_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = body_box.text_frame
        text_frame.word_wrap = True

        # 폰트 설정
        font_family = self._map_font(style.fonts.body_family)
        font_size = style.fonts.body_size[0]
        text_color = style.colors.text or style.colors.primary

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = bullet
            p.level = 0
            p.font.name = font_family
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*self._hex_to_rgb(text_color))

    def _apply_signature_elements(self, slide, style: Style):
        """시그니처 요소 적용 (간단 버전)

        추후 고도화:
        - Glassmorphism: 투명 카드
        - Swiss International: 왼쪽 레드 바
        - Neo-Brutalism: 두꺼운 테두리
        """
        # 스타일별 시그니처 요소 적용
        if style.id == "swiss_international":
            # 왼쪽 레드 바
            left = Inches(0)
            top = Inches(0)
            width = Inches(0.3)
            height = Inches(7.5)

            shape = slide.shapes.add_shape(
                1,  # Rectangle
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb("#E8000D"))
            shape.line.fill.background()

        elif style.id == "glassmorphism":
            # 배경 글로우 효과 (간단 원)
            left = Inches(7)
            top = Inches(5)
            width = Inches(2)
            height = Inches(2)

            shape = slide.shapes.add_shape(
                9,  # Oval
                left, top, width, height
            )
            shape.fill.solid()
            # 투명 효과는 python-pptx에서 제한적
            shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb("#6B21A8"))
            shape.line.fill.background()

        elif style.id == "neo_brutalism":
            # 두꺼운 검은 테두리 (슬라이드 전체)
            slide_width = Inches(10)
            slide_height = Inches(7.5)

            # python-pptx에서는 슬라이드 테두리 직접 추가 어려움
            # 대신 내부에 사각형 추가
            pass

    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """HEX 색상 → RGB 튜플

        Args:
            hex_color: #RRGGBB 또는 #RGB

        Returns:
            (R, G, B) 튜플
        """
        # 불필요한 문자 제거
        hex_color = hex_color.strip().lstrip("#")

        # #RGB → #RRGGBB
        if len(hex_color) == 3:
            hex_color = "".join([c*2 for c in hex_color])

        # 투명도 제거 (@ 이후)
        if "@" in hex_color:
            hex_color = hex_color.split("@")[0]

        # HEX → RGB
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)

        return (r, g, b)

    def _map_font(self, font_name: str) -> str:
        """폰트 매핑 (스타일 폰트 → 시스템 폰트)"""
        return self.FONT_MAPPING.get(font_name, "Calibri")

    def _get_alignment(self, alignment: str) -> int:
        """정렬 매핑"""
        alignment_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY
        }
        return alignment_map.get(alignment.lower(), PP_ALIGN.LEFT)


# 테스트 코드
if __name__ == "__main__":
    from pptx_style_recommender import PPTXStyleRecommender

    print("=" * 60)
    print("PPTX Style Renderer Test")
    print("=" * 60)

    # 1. 추천기 + 렌더러 초기화
    recommender = PPTXStyleRecommender()
    renderer = PPTXStyleRenderer()

    # 2. 테스트 케이스
    test_cases = [
        {
            "content": "구해줘 부동산 월간 실적 보고서",
            "title": "2026년 3월 월간 실적",
            "bullets": [
                "총 매출: 3,500만원 (+15% YoY)",
                "신규 계약: 7건",
                "지식산업센터 임대 상담: 12건",
                "만족도: 4.8/5.0"
            ]
        },
        {
            "content": "AI 자동화 시스템 소개",
            "title": "AI 자동화 시스템",
            "bullets": [
                "하네스 엔지니어링 기반",
                "7개 스크립트 리팩토링 완료",
                "249줄 코드 감소 (-19%)",
                "실제 운영 테스트 통과"
            ]
        }
    ]

    for i, test in enumerate(test_cases, 1):
        print(f"\n📊 Test Case {i}: {test['content']}")

        # 스타일 추천
        style = recommender.get_top_style(test['content'])
        print(f"   추천 스타일: {style.name}")

        # PPTX 생성
        prs = renderer.create_presentation(
            title=test['title'],
            content="",
            style=style,
            bullets=test['bullets']
        )

        # 저장
        output_file = f"/tmp/test_pptx_{i}_{style.id}.pptx"
        renderer.save(output_file)

        print(f"   저장: {output_file}")

    print("\n" + "=" * 60)
    print("✅ 테스트 완료")
    print("=" * 60)
